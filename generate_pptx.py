from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Hello World"
prs.save('tests/e2e/test-dummy.pptx')
